"""
Generate PPTX presentation for CV AOL project.
Run: python generate_ppt.py
Output: report/presentation.pptx
"""

import subprocess
import sys

# Auto-install python-pptx if not available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

import os

# ─── Colors ───────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT     = RGBColor(0x23, 0x6F, 0xB4)   # blue
ACCENT2    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
MED_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
HIGHLIGHT  = RGBColor(0xE8, 0xF4, 0xFD)   # light blue bg for table

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # truly blank layout


# ─── Helper: background ───────────────────────────────
def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─── Helper: add text box ─────────────────────────────
def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


# ─── Helper: accent bar top ───────────────────────────
def add_accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ─── Helper: bullet list ──────────────────────────────
def add_bullets(slide, items, left, top, width, height,
                size=20, color=BLACK, indent=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("• " if indent else "") + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


# ─── Helper: section label ────────────────────────────
def add_label(slide, text, left, top, width=Inches(3), height=Inches(0.4),
              color=ACCENT):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


# ─── Helper: divider line ─────────────────────────────
def add_divider(slide, top, color=MED_GRAY):
    line = slide.shapes.add_shape(
        1,
        Inches(0.5), top,
        Inches(12.33), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ══════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)

# Left accent panel
panel = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
panel.fill.solid()
panel.fill.fore_color.rgb = ACCENT
panel.line.fill.background()

# Top bar
add_accent_bar(sl, ACCENT2)

# Title
add_text(sl, "Automated Contrast Correction System",
         Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
         size=36, bold=True, color=ACCENT)

add_text(sl, "for Under/Over-exposed Photography\nusing Optimized Histogram Equalization",
         Inches(0.8), Inches(2.3), Inches(11), Inches(1.0),
         size=26, bold=False, color=BLACK)

add_divider(sl, Inches(3.5), ACCENT2)

add_text(sl, "COMP7116001 — Computer Vision  |  Semester Even 2025/2026",
         Inches(0.8), Inches(3.7), Inches(9), Inches(0.5),
         size=16, color=RGBColor(0x55, 0x55, 0x55))

add_text(sl, "Dosen: Fiqri Ramadhan Tambunan, S.Kom., M.Kom",
         Inches(0.8), Inches(4.2), Inches(9), Inches(0.5),
         size=16, color=RGBColor(0x55, 0x55, 0x55))

add_text(sl, "Team 2  |  Ahmad Daffa  ·  Andrey  ·  Ezra  ·  Jason  ·  Keanu  ·  Rafifdiya",
         Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
         size=16, bold=True, color=BLACK)

add_text(sl, "Binus University",
         Inches(0.8), Inches(6.5), Inches(4), Inches(0.5),
         size=14, color=RGBColor(0x99, 0x99, 0x99))


# ══════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Problem Statement", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)

add_divider(sl, Inches(1.0))

problems = [
    "Foto underexposed — detail di area shadow hilang",
    "Foto overexposed — highlight blown out, warna pucat",
    "Histogram Equalization klasik butuh parameter manual tiap gambar",
    "Tidak ada sistem yang otomatis deteksi kondisi & koreksi sesuai kebutuhan",
]
add_bullets(sl, problems, Inches(0.6), Inches(1.2), Inches(12), Inches(3.5), size=22)

add_divider(sl, Inches(5.2), ACCENT2)
add_text(sl, "[ Placeholder: 3 contoh foto — Underexposed | Normal | Overexposed ]",
         Inches(0.6), Inches(5.3), Inches(12), Inches(1.5),
         size=16, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIVES & DATASET
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Objectives & Dataset", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

add_label(sl, "TUJUAN", Inches(0.6), Inches(1.1), Inches(1.8))
objectives = [
    "Bangun pipeline koreksi kontras otomatis (GHE + CLAHE)",
    "ML Classifier: deteksi kondisi exposure secara otomatis",
    "ML Predictor: pilih parameter CLAHE optimal per gambar",
    "Deploy sebagai web app interaktif (Streamlit)",
]
add_bullets(sl, objectives, Inches(0.6), Inches(1.65), Inches(12), Inches(2.2), size=20)

add_divider(sl, Inches(4.0))

add_label(sl, "DATASET", Inches(0.6), Inches(4.15), Inches(1.8))
dataset_info = [
    "LOL (Low-Light) Dataset — Wei et al., BMVC 2018",
    "485 training pairs + 15 test pairs",
    "Format: pasangan PNG (low-light ↔ normal-light)",
]
add_bullets(sl, dataset_info, Inches(0.6), Inches(4.7), Inches(12), Inches(1.8), size=20)


# ══════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM PIPELINE
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "System Pipeline", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

steps = [
    ("Input Foto", Inches(0.4)),
    ("Ekstrak 10 Fitur Histogram (L channel)", Inches(2.0)),
    ("ML Classifier  →  Under / Over / Normal", Inches(3.6)),
    ("ML Predictor   →  clip_limit + tile_size", Inches(5.2)),
    ("CLAHE / GHE pada L channel (LAB space)", Inches(6.8)),
    ("Output + Metrics (PSNR, SSIM, Entropy)", Inches(8.4)),
]

for label, left in steps:
    # Box
    box = sl.shapes.add_shape(1, left, Inches(1.8), Inches(1.4), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = HIGHLIGHT
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.color.rgb = BLACK
    run.font.name = "Calibri"

    # Arrow (except last)
    if left != Inches(8.4):
        arrow = sl.shapes.add_shape(
            1, left + Inches(1.42), Inches(2.1),
            Inches(0.55), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT2
        arrow.line.fill.background()

add_text(sl,
         "Mode Manual: user override clip_limit & tile_size via slider kapan saja",
         Inches(0.6), Inches(3.2), Inches(12), Inches(0.5),
         size=16, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

add_text(sl, "[ Placeholder: ganti diagram panah di atas dengan flowchart visual di Canva ]",
         Inches(0.6), Inches(4.0), Inches(12), Inches(2.8),
         size=15, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 5 — GHE & CLAHE
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "CV Methods: GHE & CLAHE", Inches(0.6), Inches(0.25),
         Inches(10), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

# Two-column layout
add_label(sl, "GHE — Global", Inches(0.6), Inches(1.1), Inches(2.5))
ghe = [
    "Equalizes intensitas piksel secara global",
    "s = (L−1) × CDF(r)",
    "Cepat, mudah — dipakai sebagai baseline",
    "Kelemahan: rawan over-enhance & noise",
]
add_bullets(sl, ghe, Inches(0.6), Inches(1.65), Inches(5.7), Inches(2.5), size=18)

add_label(sl, "CLAHE — Adaptive", Inches(6.9), Inches(1.1), Inches(2.5))
clahe = [
    "Equalize per tile lokal (grid N×N)",
    "clip_limit: batas redistribusi histogram",
    "tile_size: ukuran grid (4×4, 8×8, 16×16)",
    "Lebih natural — noise terkontrol",
]
add_bullets(sl, clahe, Inches(6.9), Inches(1.65), Inches(5.7), Inches(2.5), size=18)

# Vertical divider
vdiv = sl.shapes.add_shape(1, Inches(6.6), Inches(1.1), Inches(0.02), Inches(2.8))
vdiv.fill.solid()
vdiv.fill.fore_color.rgb = MED_GRAY
vdiv.line.fill.background()

add_divider(sl, Inches(4.5), ACCENT2)
add_text(sl, "[ Placeholder: histogram comparison — Original | After GHE | After CLAHE ]",
         Inches(0.6), Inches(4.6), Inches(12), Inches(1.8),
         size=15, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 6 — EXPERIMENT RESULTS TABLE
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Experiment Results", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

add_text(sl,
         "Grid search: clip ∈ {1.0, 2.0, 3.0, 4.0}  ×  tile ∈ {4×4, 8×8, 16×16}  |  dijalankan pada 485 gambar train",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         size=15, color=RGBColor(0x55, 0x55, 0x55))

# Table
rows = 6
cols = 5
table = sl.shapes.add_table(rows, cols,
    Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.2)).table

headers = ["Method", "clip_limit", "tile_size", "PSNR (dB)", "SSIM"]
data = [
    ["GHE (baseline)", "—",   "—",      "15.33", "0.4431"],
    ["CLAHE",          "1.0", "4×4",    "8.61",  "0.2933"],
    ["CLAHE",          "2.0", "4×4",    "9.27",  "0.3659"],
    ["CLAHE",          "3.0", "4×4",    "9.91",  "0.4171"],
    ["CLAHE  ★ BEST",  "4.0", "4×4",   "10.54",  "0.4519"],
]
col_widths = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.0)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = h
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

for ri, row in enumerate(data):
    is_best = ri == 4
    for ci, val in enumerate(row):
        cell = table.cell(ri + 1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HIGHLIGHT if is_best else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        run.font.size = Pt(15)
        run.font.bold = is_best
        run.font.color.rgb = ACCENT if is_best else BLACK
        run.font.name = "Calibri"

add_text(sl,
         "CLAHE clip=4.0 tile=4×4 unggul SSIM (0.4519 > 0.4431 GHE) — SSIM lebih relevan untuk kualitas visual manusia",
         Inches(0.6), Inches(5.1), Inches(12), Inches(0.5),
         size=16, bold=True, color=ACCENT2)


# ══════════════════════════════════════════════════════
# SLIDE 7 — ML COMPONENTS
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Machine Learning Components", Inches(0.6), Inches(0.25),
         Inches(10), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

# Classifier box
clf_box = sl.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(5.9), Inches(3.2))
clf_box.fill.solid()
clf_box.fill.fore_color.rgb = HIGHLIGHT
clf_box.line.color.rgb = ACCENT

add_text(sl, "Auto Exposure CLASSIFIER",
         Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.5),
         size=18, bold=True, color=ACCENT)
clf_items = [
    "Model    : Random Forest (100 trees)",
    "Input    : 10 fitur histogram (L channel)",
    "Output   : under / over / normal",
    "Akurasi  : 96.28% ± 1.61%  ✓",
]
add_bullets(sl, clf_items, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.2),
            size=17, indent=False)

# Predictor box
prd_box = sl.shapes.add_shape(1, Inches(6.9), Inches(1.2), Inches(5.9), Inches(3.2))
prd_box.fill.solid()
prd_box.fill.fore_color.rgb = HIGHLIGHT
prd_box.line.color.rgb = ACCENT2

add_text(sl, "Auto Parameter PREDICTOR",
         Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.5),
         size=18, bold=True, color=ACCENT2)
prd_items = [
    "Model    : RF Regressor + RF Classifier",
    "Input    : 10 fitur histogram (L channel)",
    "Output   : clip_limit + tile_size",
    "MAE clip : 0.0837  |  Tile Acc : 99.59%  ✓",
]
add_bullets(sl, prd_items, Inches(7.1), Inches(1.9), Inches(5.5), Inches(2.2),
            size=17, indent=False)

add_divider(sl, Inches(4.6))
add_text(sl,
         "10 Fitur: mean · std · skewness · kurtosis · p10 · p25 · p75 · p90 · contrast ratio · entropy",
         Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
         size=15, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 8 — WEB APP DEMO
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Web Application — Demo", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

features = [
    "Upload foto atau ambil dari kamera browser",
    "Mode Auto (ML) atau Manual (slider clip + tile)",
    "Label exposure otomatis: UNDEREXPOSED / NORMAL / OVEREXPOSED",
    "Side-by-side comparison: Original | GHE | CLAHE",
    "Metrics: PSNR + SSIM (jika ada GT)  +  Uniformity + Entropy",
    "Download: best result / GHE / CLAHE",
]
add_bullets(sl, features, Inches(0.6), Inches(1.2), Inches(5.5), Inches(4.5), size=18)

# Screenshot placeholder (right side)
ph = sl.shapes.add_shape(1, Inches(6.4), Inches(1.2), Inches(6.3), Inches(5.0))
ph.fill.solid()
ph.fill.fore_color.rgb = LIGHT_GRAY
ph.line.color.rgb = MED_GRAY
tf = ph.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "[ Screenshot App ]\nStreamlit run app.py\nlalu ambil screenshot"
run.font.size = Pt(15)
run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
run.font.name = "Calibri"

add_text(sl, "Stack: Python · OpenCV · Streamlit · scikit-learn · scikit-image",
         Inches(0.6), Inches(6.1), Inches(5.5), Inches(0.4),
         size=14, color=RGBColor(0x77, 0x77, 0x77))


# ══════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)
add_accent_bar(sl)

add_text(sl, "Conclusion", Inches(0.6), Inches(0.25),
         Inches(8), Inches(0.6), size=30, bold=True, color=ACCENT)
add_divider(sl, Inches(1.0))

add_label(sl, "HASIL", Inches(0.6), Inches(1.1), Inches(1.5))
results = [
    "Pipeline GHE + CLAHE berhasil koreksi kontras otomatis",
    "Parameter optimal: CLAHE clip=4.0 tile=4×4  (SSIM = 0.4519)",
    "ML Classifier akurasi 86.74% — exposure terdeteksi otomatis",
    "ML Predictor MAE 0.08 — parameter dipilih tanpa tuning manual",
    "Web app berjalan — upload, kamera, download tersedia",
]
add_bullets(sl, results, Inches(0.6), Inches(1.65), Inches(12), Inches(2.2), size=19)

add_divider(sl, Inches(4.1))

add_label(sl, "LIMITASI & FUTURE WORK", Inches(0.6), Inches(4.2), Inches(3.2))
limits = [
    "Dataset overexposed masih synthesized — bukan real",
    "Future: deep learning (RetinexNet), real data, video support",
]
add_bullets(sl, limits, Inches(0.6), Inches(4.75), Inches(12), Inches(1.2), size=17,
            color=RGBColor(0x55, 0x55, 0x55))


# ══════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, WHITE)

# Left accent panel
panel = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
panel.fill.solid()
panel.fill.fore_color.rgb = ACCENT
panel.line.fill.background()

add_accent_bar(sl, ACCENT2)

add_text(sl, "Thank You", Inches(0.8), Inches(1.5),
         Inches(11), Inches(1.2), size=48, bold=True, color=ACCENT)

add_text(sl, "Any Questions?", Inches(0.8), Inches(2.8),
         Inches(8), Inches(0.7), size=26, color=BLACK)

add_divider(sl, Inches(3.8), ACCENT2)

add_text(sl, "GitHub  :  [ isi link repo ]",
         Inches(0.8), Inches(4.0), Inches(9), Inches(0.45), size=17,
         color=RGBColor(0x55, 0x55, 0x55))
add_text(sl, "Demo    :  [ isi link Streamlit deploy ]",
         Inches(0.8), Inches(4.5), Inches(9), Inches(0.45), size=17,
         color=RGBColor(0x55, 0x55, 0x55))

add_text(sl,
         "Ahmad Daffa  ·  Andrey  ·  Ezra  ·  Jason  ·  Keanu  ·  Rafifdiya\nTeam 2  |  COMP7116001  |  Binus University",
         Inches(0.8), Inches(5.4), Inches(11), Inches(1.0),
         size=16, color=RGBColor(0x77, 0x77, 0x77))


# ─── Save ─────────────────────────────────────────────
out_path = os.path.join("report", "presentation.pptx")
os.makedirs("report", exist_ok=True)
prs.save(out_path)
print(f"Saved: {out_path}")
print("Done! Open report/presentation.pptx")
